import os, json, sqlite3, secrets, threading, webbrowser, base64, io, logging, random
from datetime import datetime, date
from flask import Flask, render_template, request, jsonify, redirect, url_for, session, send_file, Response
import bcrypt  # Seguridad robusta
from faker import Faker # Datos ficticios

fake = Faker('es_ES')

app = Flask(__name__, template_folder='templates', static_folder='static')

# --- CONFIGURACIÓN DE SEGURIDAD Y RUTAS ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH  = os.path.join(BASE_DIR, 'datos', 'semillero.db')

# Secret Key Persistente
_SECRET_FILE = os.path.join(BASE_DIR, 'datos', '.secret_key')
if os.path.exists(_SECRET_FILE):
    with open(_SECRET_FILE, 'r') as _f: 
        app.secret_key = _f.read().strip()
else:
    os.makedirs(os.path.join(BASE_DIR, 'datos'), exist_ok=True)
    _key = secrets.token_hex(32)
    with open(_SECRET_FILE, 'w') as _f: 
        _f.write(_key)
    app.secret_key = _key

@app.after_request
def add_header(response):
    response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, public, max-age=0'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    return response

# Logging Centralizado
_LOG_DIR = os.path.join(BASE_DIR, 'datos')
os.makedirs(_LOG_DIR, exist_ok=True)
logging.basicConfig(
    filename=os.path.join(_LOG_DIR, 'app.log'),
    level=logging.ERROR,
    format='%(asctime)s [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger('semillero')

# --- DATOS GEOGRÁFICOS ---
MUNICIPIOS_COORDS = {
    "Alberto Adriani": [8.62, -71.65], "Libertador": [8.6, -71.15], "Iribarren": [10.07, -69.32],
    "San Cristóbal": [7.77, -72.22], "Maracaibo": [10.63, -71.65], "Valencia": [10.18, -68.0],
    "Libertador (Caracas)": [10.48, -66.9], "Barcelona": [10.13, -64.68], "Maturín": [9.75, -63.18],
    "Guanare": [9.04, -69.74], "Barinas": [8.62, -70.21], "Caroní": [8.3, -62.73],
    "Tucupita": [9.06, -62.05], "Cumaná": [10.46, -64.17], "Porlamar": [10.96, -63.85],
    "San Fernando": [7.9, -67.47], "Crespo": [10.07, -69.32], "Palavecino": [10.18, -69.05],
    "Girardot": [10.24, -67.59], "Sucre (Miranda)": [10.5, -66.8],
}

MUNICIPIOS_VE = {
    "Amazonas": ["Alto Orinoco", "Atabapo", "Atures", "Autana", "Manapiare", "Maroa", "Rio Negro"],
    "Anzoátegui": ["Anaco", "Barcelona", "Bergantín", "Bolívar", "Cajigal", "Diego Bautista Urbaneja", "Freites", "Guanipa", "Guanta", "Independencia", "Libertad", "Miranda", "Monagas", "Páez", "Peñalver", "Píritu", "Simón Bolívar", "Sotillo"],
    "Apure": ["Achaguas", "Biruaca", "Muñoz", "Páez", "Pedro Camejo", "Pulido", "Rómulo Gallegos", "San Fernando"],
    "Aragua": ["Bolívar", "Camatagua", "Francisco Linares Alcántara", "Girardot", "José Ángel Lamas", "José Félix Ribas", "Libertador", "Mario Briceño Iragorry", "Ocumare de la Costa", "San Casimiro", "San Sebastián", "Santiago Mariño", "Sucre", "Tovar", "Urdaneta", "Zamora"],
    "Barinas": ["Alberto Arvelo Torrealba", "Arismendi", "Barinas", "Bolívar", "Cruz Paredes", "Ezequiel Zamora", "Obispos", "Pedraza", "Rojas", "Sosa", "Sucre"],
    "Bolívar": ["Angostura", "Caroní", "Cedeño", "El Callao", "Gran Sabana", "Heres", "Independencia", "Padre Pedro Chien", "Piar", "Raúl Leoni", "Roscio", "Sifontes", "Sucre"],
    "Carabobo": ["Bejuma", "Carlos Arvelo", "Diego Ibarra", "Guacara", "Juan José Mora", "Libertador", "Los Guayos", "Miranda", "Montalbán", "Naguanagua", "Puerto Cabello", "San Diego", "San Joaquín", "Valencia"],
    "Cojedes": ["Anzoátegui", "Ezequiel Zamora", "Girardot", "Lima Blanco", "Pao de San Juan Bautista", "Ricaurte", "Rómulo Gallegos", "San Carlos", "Tinaco"],
    "Delta Amacuro": ["Antonio Díaz", "Casacoima", "Pedernales", "Tucupita"],
    "Distrito Capital": ["Libertador"],
    "Falcón": ["Acosta", "Bolívar", "Buchivacoa", "Carirubana", "Colina", "Dabajuro", "Democracia", "Falcón", "Federación", "Jacura", "Los Taques", "Mauroa", "Miranda", "Palmasola", "Petit", "Piritu", "San Francisco", "Sucre", "Tocópero", "Unión", "Urumaco", "Zamora"],
    "Guárico": ["Camaguán", "Chaguaramas", "El Socorro", "Francisco de Miranda", "Julián Mellado", "Las Mercedes", "Leonardo Infante", "Monagas", "Ortiz", "Pedro Zaraza", "San Gerónimo de Guayabal", "San José de Guaribe", "Santa María de Ipire", "Urdaneta"],
    "Lara": ["Andrés Eloy Blanco", "Crespo", "Iribarren", "Jiménez", "Morán", "Palavecino", "Simón Planas", "Torres", "Urdaneta"],
    "Mérida": ["Alberto Adriani", "Andrés Bello", "Aricagua", "Arzobispo Chacón", "Campo Elías", "Caracciolo Parra Olmedo", "Cardenal Quintero", "Guaraque", "Julio César Salas", "Justo Briceño", "Libertador", "Miranda", "Obispo Ramos de Lora", "Padre Noguera", "Pueblo Llano", "Rangel", "Rivas Dávila", "Santos Marquina", "Sucre", "Tovar", "Tulio Febres Cordero", "Zea"],
    "Miranda": ["Acevedo", "Andrés Bello", "Baruta", "Brión", "Buroz", "Carrizal", "Chacao", "Cristóbal Rojas", "El Hatillo", "Guaicaipuro", "Independencia", "Lander", "Los Salias", "Páez", "Paz Castillo", "Pedro Gual", "Plaza", "Simón Bolívar", "Sucre", "Urdaneta", "Zamora"],
    "Monagas": ["Acosta", "Aguasay", "Bolívar", "Caripe", "Cedeño", "Ezequiel Zamora", "Libertador", "Maturín", "Piar", "Punceres", "Santa Bárbara", "Sotillo", "Uracoa"],
    "Nueva Esparta": ["Arismendi", "Díaz", "García", "Gómez", "Gual", "Maneiro", "Marcano", "Mariño", "Morales", "Península de Macanao", "Tubores", "Villalba"],
    "Portuguesa": ["Araure", "Esteller", "Guanare", "Guanarito", "Monseñor José Vicente de Unda", "Ospino", "Páez", "Papelón", "San Genaro de Boconoíto", "San Rafael de Onoto", "Santa Rosalía", "Sucre", "Turen"],
    "Sucre": ["Andrés Mata", "Arismendi", "Benítez", "Bermúdez", "Bolívar", "Cruz Salmerón Acosta", "Cumanacoa", "Libertador", "Mariño", "Mejías", "Montes", "Ribero", "Sucre", "Valdez"],
    "Táchira": ["Andrés Bello", "Ayacucho", "Bolívar", "Cárdenas", "Córdoba", "Fernández Feo", "Francisco de Miranda", "García de Hevia", "Guásimos", "Independencia", "Jáuregui", "José María Vargas", "Junín", "Libertad", "Libertador", "Lobatera", "Michelena", "Panamericano", "Pedro María Ureña", "Rafael Urdaneta", "Ribas", "San Cristóbal", "Simón Rodríguez", "Sucre", "Torbes", "Uribante", "Zerpa"],
    "Trujillo": ["Andrés Bello", "Boconó", "Candelaria", "Carache", "Escuque", "La Ceiba", "Miranda", "Motatán", "Pampán", "Pampanito", "Rafael Rangel", "San Rafael de Carvajal", "Sucre", "Trujillo", "Urdaneta", "Valera"],
    "Vargas": ["Vargas"],
    "Yaracuy": ["Bolívar", "Bruzual", "Cocorote", "Independencia", "La Trinidad", "Manuel Monge", "Nirgua", "Páez", "San Felipe", "Sucre", "Urachiche", "Veroes"],
    "Zulia": ["Almirante Padilla", "Baralt", "Cabimas", "Catatumbo", "Colón", "Democracia", "Francisco Javier Pulgar", "Guajira", "Jesús Enrique Lossada", "Jesús María Semprún", "La Cañada de Urdaneta", "Lagunillas", "Machiques de Perijá", "Mara", "Maracaibo", "Miranda", "Páez", "Rosario de Perijá", "San Francisco", "Santa Rita", "Simón Bolívar", "Sucre", "Valmore Rodríguez"],
}

# --- FUNCIONES DE SEGURIDAD Y DB ---

def hash_pw(pw):
    return bcrypt.hashpw(pw.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def check_pw(pw, hashed):
    try:
        return bcrypt.checkpw(pw.encode('utf-8'), hashed.encode('utf-8'))
    except Exception:
        return False

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def get_cfg(k):
    conn = get_db()
    r = conn.execute("SELECT valor FROM config_sistema WHERE clave=?", (k,)).fetchone()
    conn.close()
    return r['valor'] if r else None

def is_setup(): 
    return get_cfg('configurado') == '1'

def login_required(f):
    from functools import wraps
    @wraps(f)
    def d(*a, **k):
        if not session.get('usuario'): 
            return redirect(url_for('login'))
        return f(*a, **k)
    return d

# --- HELPERS DE VALIDACIÓN ---

def validar_password(pw):
    if not pw: return 'La contraseña no puede estar vacía.'
    if len(pw) < 6: return 'La contraseña debe tener al menos 6 caracteres.'
    return None

def validar_username(u):
    import re
    if not u: return 'El nombre de usuario no puede estar vacío.'
    if len(u) < 3: return 'El usuario debe tener al menos 3 caracteres.'
    if not re.match(r'^[a-zA-Z0-9_.]+$', u):
        return 'El usuario solo puede contener letras, números, puntos y guiones bajos.'
    return None

# --- INICIALIZACIÓN DE DB Y DATOS FICTICIOS ---

def init_db():
    os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
    conn = get_db()
    conn.executescript("""
    CREATE TABLE IF NOT EXISTS config_sistema(clave TEXT PRIMARY KEY, valor TEXT);
    CREATE TABLE IF NOT EXISTS usuarios(id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL, password_hash TEXT NOT NULL,
        pregunta_seguridad TEXT, respuesta_hash TEXT, nombre_completo TEXT,
        rol TEXT DEFAULT 'admin', activo INTEGER DEFAULT 1);
    CREATE TABLE IF NOT EXISTS voceras(id INTEGER PRIMARY KEY AUTOINCREMENT,
        nombres TEXT NOT NULL, apellidos TEXT NOT NULL, cedula TEXT NOT NULL,
        cargo TEXT DEFAULT 'Vocera/o', activo INTEGER DEFAULT 1, orden INTEGER DEFAULT 1);
    CREATE TABLE IF NOT EXISTS beneficiarios(id INTEGER PRIMARY KEY AUTOINCREMENT,
        contrato TEXT UNIQUE NOT NULL, nombres TEXT NOT NULL, apellidos TEXT NOT NULL,
        cedula TEXT NOT NULL, rif TEXT NOT NULL, direccion TEXT, direccion_real TEXT,
        sector TEXT, lat REAL DEFAULT 8.623, lng REAL DEFAULT -71.651,
        estado_civil TEXT, fecha_nacimiento TEXT, actividad TEXT,
        fecha_desembolso TEXT, monto REAL DEFAULT 400, cuotas_total INTEGER DEFAULT 11,
        cuotas_pagadas INTEGER DEFAULT 0, rif_vencimiento TEXT,
        foto_cedula TEXT, foto_rif TEXT, foto_perfil TEXT, historia TEXT,
        fecha_registro TEXT DEFAULT CURRENT_TIMESTAMP, activo INTEGER DEFAULT 1);
    CREATE TABLE IF NOT EXISTS pagos(id INTEGER PRIMARY KEY AUTOINCREMENT,
        beneficiario_id INTEGER NOT NULL, contrato TEXT NOT NULL,
        numero_cuota INTEGER NOT NULL, monto_usd REAL DEFAULT 40,
        tasa_bcv REAL, monto_bs REAL, fecha_pago TEXT,
        numero_operacion TEXT, banco_origen TEXT, foto_capture TEXT,
        registrado_por TEXT DEFAULT 'Sistema',
        fecha_registro TEXT DEFAULT CURRENT_TIMESTAMP);
    CREATE TABLE IF NOT EXISTS fotos_historia(id INTEGER PRIMARY KEY AUTOINCREMENT,
        beneficiario_id INTEGER NOT NULL, imagen_b64 TEXT NOT NULL,
        descripcion TEXT, fecha TEXT, tipo TEXT DEFAULT 'historia',
        fecha_registro TEXT DEFAULT CURRENT_TIMESTAMP);
    CREATE TABLE IF NOT EXISTS movimientos_banco(id INTEGER PRIMARY KEY AUTOINCREMENT,
        tipo TEXT NOT NULL, concepto TEXT NOT NULL, monto_usd REAL NOT NULL,
        fecha TEXT NOT NULL, referencia TEXT, beneficiario_id INTEGER, notas TEXT,
        fecha_registro TEXT DEFAULT CURRENT_TIMESTAMP);
    CREATE TABLE IF NOT EXISTS contratos_generados(id INTEGER PRIMARY KEY AUTOINCREMENT,
        beneficiario_id INTEGER, contrato TEXT, ruta_archivo TEXT,
        fecha_generacion TEXT DEFAULT CURRENT_TIMESTAMP);
    CREATE TABLE IF NOT EXISTS transacciones_banco(id INTEGER PRIMARY KEY AUTOINCREMENT,
        fecha TEXT NOT NULL, referencia TEXT UNIQUE NOT NULL,
        descripcion TEXT, codigo TEXT, debito REAL DEFAULT 0.0,
        credito REAL DEFAULT 0.0, saldo REAL, reconciliado INTEGER DEFAULT 0,
        pago_id INTEGER, fecha_registro TEXT DEFAULT CURRENT_TIMESTAMP);
    """)
    # Migraciones seguras para actualizar pagos con campos de conciliacion
    try:
        conn.execute("ALTER TABLE pagos ADD COLUMN conciliado INTEGER DEFAULT 0")
    except Exception:
        pass
    try:
        conn.execute("ALTER TABLE pagos ADD COLUMN banco_id_referencia TEXT")
    except Exception:
        pass
    conn.commit()
    conn.close()

def seed_beneficiarios_ficticios():
    """Genera datos DEMO ficticios usando Faker."""
    conn = get_db()
    c = conn.cursor()
    
    # Generar 10 beneficiarios ficticios
    for i in range(1, 11):
        contrato = f"BCCPPO-{i:02d}-2026"
        perfil = fake.profile()
        nombres = perfil['name'].split()[0]
        apellidos = ' '.join(perfil['name'].split()[1:])
        cedula = f"V-{fake.random_number(digits=8, fix_len=True)}"
        rif = f"V{cedula.split('-')[1]}0"
        direccion = fake.address().replace('\n', ', ')
        sector = fake.city()
        lat = 8.62 + (random.random() * 0.05 - 0.025)
        lng = -71.65 + (random.random() * 0.05 - 0.025)
        estado_civil = random.choice(["SOLTERO", "CASADO", "DIVORCIADO", "VIUDO"])
        fecha_nac = perfil['birthdate'].strftime("%d/%m/%Y") if perfil['birthdate'] else "01/01/1980"
        
        try:
            c.execute("""INSERT OR IGNORE INTO beneficiarios
            (contrato,nombres,apellidos,cedula,rif,direccion,sector,lat,lng,
            estado_civil,fecha_nacimiento,fecha_desembolso,rif_vencimiento)
            VALUES(?,?,?,?,?,?,?,?,?,?,?,date('now'), '31/12/2028')""",
            (contrato, nombres, apellidos, cedula, rif, direccion, sector, lat, lng, estado_civil, fecha_nac))
            
            row = c.execute("SELECT id FROM beneficiarios WHERE contrato=?", (contrato,)).fetchone()
            if row:
                c.execute("""INSERT OR IGNORE INTO movimientos_banco(tipo,concepto,monto_usd,fecha,referencia,beneficiario_id)
                SELECT 'EGRESO',?,400,date('now'),?,? WHERE NOT EXISTS
                (SELECT 1 FROM movimientos_banco WHERE referencia=? AND tipo='EGRESO')""",
                (f"Desembolso {contrato}", contrato, row['id'], contrato))
        except Exception as e:
            logger.error(f"Error generando dato ficticio {i}: {e}")
            continue
            
    conn.commit()
    conn.close()

# --- RUTAS DE AUTENTICACIÓN ---

@app.route('/')
@app.route('/index')
@login_required
def index(): 
    return render_template('index_v23.html')

@app.route('/setup', methods=['GET','POST'])
def setup():
    if is_setup(): 
        return redirect(url_for('login'))
    
    errores = {}
    form = {}
    
    if request.method == 'POST':
        d = request.form
        form = d
        
        if not d.get('nombre_comuna','').strip(): errores['nombre_comuna'] = 'Obligatorio.'
        if not d.get('nombre_banco','').strip(): errores['nombre_banco'] = 'Obligatorio.'
        if not d.get('estado','').strip(): errores['estado'] = 'Seleccione estado.'
        if not d.get('municipio','').strip(): errores['municipio'] = 'Seleccione municipio.'
        if not d.get('ciudad','').strip(): errores['ciudad'] = 'Obligatorio.'
        
        err_u = validar_username(d.get('usuario','').strip())
        if err_u: errores['usuario'] = err_u
        
        err_p = validar_password(d.get('password',''))
        if err_p: errores['password'] = err_p
        
        if not d.get('pregunta','').strip(): errores['pregunta'] = 'Obligatorio.'
        if not d.get('respuesta', '').strip(): errores['respuesta'] = 'Obligatorio.'

        if not errores:
            conn = get_db()
            c = conn.cursor()
            
            for k,v in {
                'configurado':'1',
                'nombre_comuna': d.get('nombre_comuna','').strip(),
                'estado': d.get('estado',''),
                'municipio': d.get('municipio',''),
                'ciudad': d.get('ciudad','').strip(),
                'api_key': d.get('api_key',''),
                'nombre_banco': d.get('nombre_banco','').strip()
            }.items():
                c.execute("INSERT OR REPLACE INTO config_sistema(clave,valor) VALUES(?,?)",(k,v))
            
            c.execute("""INSERT OR REPLACE INTO usuarios
               (username,password_hash,pregunta_seguridad,respuesta_hash,nombre_completo,rol)
               VALUES(?,?,?,?,?,?)""",
                (d.get('usuario','').strip(),
                 hash_pw(d.get('password','')),
                 d.get('pregunta'),
                 hash_pw(d.get('respuesta','').lower().strip()),
                 d.get('nombre_completo','').strip(),
                 'admin')
            )
            
            # Validar que el usuario se creó correctamente
            verificar_usuario = c.execute("SELECT id FROM usuarios WHERE username=?", 
                                         (d.get('usuario','').strip(),)).fetchone()
            if not verificar_usuario:
                conn.rollback()
                conn.close()
                errores['general'] = 'Error al crear el usuario. Por favor, intenta nuevamente.'
                return render_template('setup.html', errores=errores, valores=d, 
                                     estados=list(MUNICIPIOS_VE.keys()))
            
            for i in range(1,4):
                nom = d.get(f'vocera{i}_nombres','').strip()
                ape = d.get(f'vocera{i}_apellidos','').strip()
                ci  = d.get(f'vocera{i}_cedula','').strip()
                if nom and ci:
                    c.execute("INSERT INTO voceras(nombres,apellidos,cedula,orden) VALUES(?,?,?,?)",(nom,ape,ci,i))
            
            conn.commit()
            conn.close()
            
            seed_beneficiarios_ficticios()
            
            return redirect(url_for('login'))

    return render_template('setup.html',
        estados=list(MUNICIPIOS_VE.keys()),
        municipios_json=json.dumps(MUNICIPIOS_VE),
        errores=errores,
        form=form)

@app.route('/login', methods=['GET','POST'])
def login():
    if not is_setup(): 
        return redirect(url_for('setup'))
    error = None
    username_val = ''
    
    if request.method == 'POST':
        username_val = request.form.get('username','').strip()
        password     = request.form.get('password','')
        
        if not username_val or not password:
            error = 'Debes ingresar usuario y contraseña.'
        else:
            conn = get_db()
            u = conn.execute("SELECT * FROM usuarios WHERE username=? AND activo=1",(username_val,)).fetchone()
            conn.close()
            
            if u and check_pw(password, u['password_hash']):
                session['usuario'] = u['username']
                session['nombre']  = u['nombre_completo'] or u['username']
                session['rol']     = u['rol']
                return redirect(url_for('index'))
            error = 'Usuario o contraseña incorrectos.'
            logger.warning(f"Intento de login fallido para usuario: {username_val}")

    return render_template('login.html',
        error=error,
        username_val=username_val,
        nombre_comuna=get_cfg('nombre_comuna') or 'SemilleroComunal')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

@app.route('/recuperar', methods=['GET','POST'])
def recuperar():
    if request.method == 'POST':
        paso     = request.form.get('paso','1')
        username = request.form.get('username','').strip()
        
        if paso == '1':
            if not username:
                return render_template('recuperar.html', paso=1, error='Ingresa tu usuario.')
            conn = get_db()
            u = conn.execute("SELECT pregunta_seguridad FROM usuarios WHERE username=? AND activo=1",(username,)).fetchone()
            conn.close()
            if u:
                return render_template('recuperar.html', paso=2, username=username, pregunta=u['pregunta_seguridad'])
            return render_template('recuperar.html', paso=1, error='Usuario no encontrado.')

        elif paso == '2':
            respuesta    = request.form.get('respuesta','').lower().strip()
            nueva_pw     = request.form.get('nueva_password','')
            confirm_pw   = request.form.get('confirmar_password','')

            conn = get_db()
            u = conn.execute("SELECT * FROM usuarios WHERE username=? AND activo=1",(username,)).fetchone()

            if not u:
                conn.close()
                return render_template('recuperar.html', paso=1, error='Usuario inválido.')

            if not check_pw(respuesta, u['respuesta_hash']):
                conn.close()
                return render_template('recuperar.html', paso=2, username=username,
                     pregunta=u['pregunta_seguridad'], error='Respuesta incorrecta.')

            err_p = validar_password(nueva_pw)
            if err_p:
                conn.close()
                return render_template('recuperar.html', paso=2, username=username,
                    pregunta=u['pregunta_seguridad'], error=err_p)

            if nueva_pw != confirm_pw:
                conn.close()
                return render_template('recuperar.html', paso=2, username=username,
                    pregunta=u['pregunta_seguridad'], error='Las contraseñas no coinciden.')

            conn.execute("UPDATE usuarios SET password_hash=? WHERE username=?", (hash_pw(nueva_pw), username))
            conn.commit()
            conn.close()
            return render_template('recuperar.html', paso=3, mensaje='¡Contraseña actualizada!')

    return render_template('recuperar.html', paso=1)

# --- APIs DE CONFIGURACIÓN Y USUARIOS ---

@app.route('/api/config')
@login_required
def api_config():
    conn=get_db()
    rows=conn.execute("SELECT clave,valor FROM config_sistema").fetchall()
    conn.close()
    cfg={r['clave']:r['valor'] for r in rows}
    mun=cfg.get('municipio','')
    coords=MUNICIPIOS_COORDS.get(mun,[8.623,-71.651])
    cfg['map_lat']=coords[0]
    cfg['map_lng']=coords[1]
    cfg['usuario']=session.get('nombre', session.get('usuario',''))
    cfg['has_api_key'] = bool(cfg.get('api_key'))
    cfg['has_groq_key'] = bool(cfg.get('groq_api_key'))
    cfg.pop('api_key', None)
    cfg.pop('groq_api_key', None)
    return jsonify(cfg)

@app.route('/api/municipios/<estado>')
def api_municipios(estado): 
    return jsonify(MUNICIPIOS_VE.get(estado,[]))

@app.route('/api/voceras')
@login_required
def api_voceras():
    conn=get_db()
    v=conn.execute("SELECT * FROM voceras WHERE activo=1 ORDER BY orden").fetchall()
    conn.close()
    return jsonify([dict(x) for x in v])

@app.route('/api/voceras/guardar', methods=['POST'])
@login_required
def guardar_voceras():
    voceras=request.json.get('voceras',[])
    conn=get_db()
    c=conn.cursor()
    c.execute("DELETE FROM voceras")
    for i,v in enumerate(voceras,1):
        if v.get('nombres') and v.get('cedula'):
            c.execute("INSERT INTO voceras(nombres,apellidos,cedula,cargo,orden) VALUES(?,?,?,?,?)",
            (v['nombres'],v.get('apellidos',''),v['cedula'],v.get('cargo','Vocera/o'),i))
    conn.commit()
    conn.close()
    return jsonify({"ok":True})

@app.route('/api/config/actualizar', methods=['POST'])
@login_required
def actualizar_config():
    d=request.json
    conn=get_db()
    c=conn.cursor()
    campos=['nombre_comuna','nombre_banco','estado','municipio','ciudad','api_preferida']
    for campo in campos:
        if campo in d:
            c.execute("INSERT OR REPLACE INTO config_sistema(clave,valor) VALUES(?,?)",(campo,d[campo]))
    if 'api_key' in d:
        c.execute("INSERT OR REPLACE INTO config_sistema(clave,valor) VALUES('api_key',?)",(d['api_key'],))
    if 'groq_api_key' in d:
        c.execute("INSERT OR REPLACE INTO config_sistema(clave,valor) VALUES('groq_api_key',?)",(d['groq_api_key'],))
    conn.commit()
    conn.close()
    return jsonify({"ok":True})

@app.route('/api/usuarios')
@login_required
def api_usuarios():
    conn=get_db()
    u=conn.execute("SELECT id,username,nombre_completo,rol,activo FROM usuarios").fetchall()
    conn.close()
    return jsonify([dict(x) for x in u])

@app.route('/api/usuarios/crear', methods=['POST'])
@login_required
def crear_usuario():
    d = request.json or {}
    username = (d.get('username') or '').strip()
    password = d.get('password','')
    nombre   = (d.get('nombre_completo') or '').strip()
    rol      = d.get('rol','operador')
    
    err_u = validar_username(username)
    if err_u: return jsonify({"error": err_u}), 400
    err_p = validar_password(password)
    if err_p: return jsonify({"error": err_p}), 400
    if rol not in ('admin','operador'):
        return jsonify({"error": "Rol inválido."}), 400

    conn = get_db()
    existe = conn.execute("SELECT id FROM usuarios WHERE username=?",(username,)).fetchone()
    if existe:
        conn.close()
        return jsonify({"error": "El usuario ya existe."}), 400

    conn.execute("INSERT INTO usuarios(username,password_hash,nombre_completo,rol,activo) VALUES(?,?,?,?,1)",
        (username, hash_pw(password), nombre, rol))
    conn.commit()
    conn.close()
    return jsonify({"ok": True, "mensaje": f"Usuario '{username}' creado."})

@app.route('/api/usuarios/<int:uid>/cambiar_password', methods=['POST'])
@login_required
def cambiar_password(uid):
    d = request.json or {}
    nueva_pw = d.get('nueva_password','')
    confirmar = d.get('confirmar_password','')
    
    err_p = validar_password(nueva_pw)
    if err_p: return jsonify({"error": err_p}), 400
    if confirmar and nueva_pw != confirmar:
        return jsonify({"error": "Las contraseñas no coinciden."}), 400

    conn = get_db()
    u = conn.execute("SELECT id FROM usuarios WHERE id=?",(uid,)).fetchone()
    if not u:
        conn.close()
        return jsonify({"error": "Usuario no encontrado."}), 404

    conn.execute("UPDATE usuarios SET password_hash=? WHERE id=?", (hash_pw(nueva_pw), uid))
    conn.commit()
    conn.close()
    return jsonify({"ok": True, "mensaje": "Contraseña actualizada."})

@app.route('/api/usuarios/<int:uid>/toggle', methods=['POST'])
@login_required
def toggle_usuario(uid):
    conn = get_db()
    u = conn.execute("SELECT activo, rol FROM usuarios WHERE id=?",(uid,)).fetchone()
    if not u:
        conn.close()
        return jsonify({"error": "Usuario no encontrado."}), 404
    
    if u['activo'] == 1 and u['rol'] == 'admin':
        admins_activos = conn.execute("SELECT COUNT(*) FROM usuarios WHERE rol='admin' AND activo=1").fetchone()[0]
        if admins_activos <= 1:
            conn.close()
            return jsonify({"error": "No puedes desactivar el único administrador."}), 400
            
    nuevo_estado = 0 if u['activo'] else 1
    conn.execute("UPDATE usuarios SET activo=? WHERE id=?", (nuevo_estado, uid))
    conn.commit()
    conn.close()
    return jsonify({"ok": True, "activo": nuevo_estado})

# --- GENERACIÓN DE DOCUMENTOS (RECIBOS) ---

@app.route('/api/recibos/generar', methods=['POST'])
@login_required
def generar_recibo():
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    
    d=request.json
    bid=d.get('beneficiario_id')
    conn=get_db()
    b=conn.execute("SELECT * FROM beneficiarios WHERE id=?",(bid,)).fetchone()
    voceras=conn.execute("SELECT * FROM voceras WHERE activo=1 ORDER BY orden LIMIT 1").fetchall()
    cfg={r['clave']:r['valor'] for r in conn.execute("SELECT clave,valor FROM config_sistema").fetchall()}
    conn.close()
    
    if not b: return jsonify({"error": "Beneficiario no encontrado"}),404

    nombre_banco=cfg.get('nombre_banco','BANCO DE LA COMUNA')
    ciudad=cfg.get('ciudad',cfg.get('municipio',''))
    estado=cfg.get('estado','')
    nombre_completo=f"{b['nombres']} {b['apellidos']}"
    primera_vocera=voceras[0] if voceras else None
    vocera_nombre=f"{primera_vocera['nombres']} {primera_vocera['apellidos']}" if primera_vocera else "______________________"
    vocera_ci=primera_vocera['cedula'] if primera_vocera else "___________"
    mes_firma=d.get('mes_firma','marzo')
    fecha_desembolso=d.get('fecha_desembolso',b['fecha_desembolso'] or f'_____ de {mes_firma} de {date.today().year}')

    doc=Document()
    for section in doc.sections:
        section.top_margin=Inches(0.6); section.bottom_margin=Inches(0.6)
        section.left_margin=Inches(1.0); section.right_margin=Inches(1.0)

    htable=doc.add_table(rows=1,cols=2)
    htable.style='Table Grid'
    for cell in htable.rows[0].cells:
        tc=cell._tc; tcPr=tc.get_or_add_tcPr()
        tcBorders=OxmlElement('w:tcBorders')
        for side in ['top','left','bottom','right']:
            border=OxmlElement(f'w:{side}')
            border.set(qn('w:val'),'none'); border.set(qn('w:sz'),'0')
            tcBorders.append(border)
        tcPr.append(tcBorders)

    lc=htable.cell(0,0)
    lc.width=Inches(1.2)
    lp=lc.paragraphs[0]; lp.alignment=WD_ALIGN_PARAGRAPH.CENTER
    logo_path=os.path.join(BASE_DIR,'static','img','semillero.png')
    if os.path.exists(logo_path):
        lp.add_run().add_picture(logo_path,width=Inches(0.9))

    tc2=htable.cell(0,1)
    tp=tc2.paragraphs[0]; tp.alignment=WD_ALIGN_PARAGRAPH.CENTER
    tr=tp.add_run(nombre_banco.upper()); tr.bold=True; tr.font.size=Pt(11)
    tp2=tc2.add_paragraph(f"{ciudad}, Estado {estado}")
    tp2.alignment=WD_ALIGN_PARAGRAPH.CENTER; tp2.runs[0].font.size=Pt(9)

    doc.add_paragraph()
    tp3=doc.add_paragraph(); tp3.alignment=WD_ALIGN_PARAGRAPH.CENTER
    tr3=tp3.add_run('CONSTANCIA DE RECEPCIÓN DE DESEMBOLSO')
    tr3.bold=True; tr3.font.size=Pt(14)
    
    doc.add_paragraph()
    intro=doc.add_paragraph()
    intro.alignment=WD_ALIGN_PARAGRAPH.JUSTIFY
    runs=[
        ('Yo, ',False), (nombre_completo,True), (', titular de la C.I. Nº ',False),
        (b['cedula'],True), (', declaro haber recibido del ',False), (nombre_banco.upper(),True),
        (' el monto del contrato ',False), (b['contrato'],True), ('.',False),
    ]
    for text,bold in runs:
        r=intro.add_run(text); r.bold=bold; r.font.size=Pt(10)

    fields=[
        ('Monto en USD:','CUATROCIENTOS DÓLARES (USD $400,00)'),
        ('Fecha de Desembolso:',fecha_desembolso),
        ('Modalidad de pago:','Transferencia / Pago Móvil'),
    ]
    ft=doc.add_table(rows=len(fields),cols=2)
    ft.style='Table Grid'
    for i,(label,value) in enumerate(fields):
        lc=ft.cell(i,0); vc=ft.cell(i,1)
        lp=lc.paragraphs[0]; lp.clear(); lr=lp.add_run(label); lr.bold=True; lr.font.size=Pt(9.5)
        vp=vc.paragraphs[0]; vp.clear(); vr=vp.add_run(value); vr.font.size=Pt(9.5)

    doc.add_paragraph()
    legal=doc.add_paragraph()
    legal.alignment=WD_ALIGN_PARAGRAPH.JUSTIFY
    lr1=legal.add_run('El beneficiario declara que los recursos serán destinados a actividades productivas comunales.')
    lr1.font.size=Pt(9)

    doc.add_paragraph()
    stbl=doc.add_table(rows=2,cols=2)
    stbl.style='Table Grid'
    for row in stbl.rows:
        for cell in row.cells:
            tc=cell._tc; tcPr=tc.get_or_add_tcPr()
            bdr=OxmlElement('w:tcBorders')
            for side in ['top','left','bottom','right']:
                b3=OxmlElement(f'w:{side}'); b3.set(qn('w:val'),'none'); b3.set(qn('w:sz'),'0'); bdr.append(b3)
            tcPr.append(bdr)
    
    p=stbl.cell(0,0).paragraphs[0]; p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    p.add_run('________________________________').font.size=Pt(11)
    bc=stbl.cell(1,0)
    par=bc.paragraphs[0]; par.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=par.add_run(nombre_completo); r.bold=True; r.font.size=Pt(9.5)
    bc.add_paragraph(f"C.I. {b['cedula']}").alignment=WD_ALIGN_PARAGRAPH.CENTER
    bc.add_paragraph("EL BENEFICIARIO").alignment=WD_ALIGN_PARAGRAPH.CENTER

    p2=stbl.cell(0,1).paragraphs[0]; p2.alignment=WD_ALIGN_PARAGRAPH.CENTER
    p2.add_run('________________________________').font.size=Pt(11)
    bc2=stbl.cell(1,1)
    par2=bc2.paragraphs[0]; par2.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r2=par2.add_run(vocera_nombre); r2.bold=True; r2.font.size=Pt(9.5)
    bc2.add_paragraph(f"C.I. {vocera_ci}").alignment=WD_ALIGN_PARAGRAPH.CENTER
    bc2.add_paragraph("POR EL BANCO").alignment=WD_ALIGN_PARAGRAPH.CENTER

    fname=f"Recibo_{b['apellidos'].replace(' ','_')}_{b['contrato']}.docx"
    fpath=os.path.join(BASE_DIR,'documentos',fname)
    doc.save(fpath)
    return send_file(fpath,as_attachment=True,download_name=fname, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

# --- DASHBOARD Y ESTADÍSTICAS ---

@app.route('/api/dashboard')
@login_required
def api_dashboard():
    conn=get_db(); c=conn.cursor()
    total=c.execute("SELECT COUNT(*) FROM beneficiarios WHERE activo=1").fetchone()[0]
    cartera=c.execute("SELECT SUM(monto) FROM beneficiarios WHERE activo=1").fetchone()[0] or 0
    cobrado=c.execute("SELECT SUM(monto_usd) FROM pagos").fetchone()[0] or 0
    mora=c.execute("SELECT COUNT(*) FROM beneficiarios WHERE activo=1 AND cuotas_pagadas=0").fetchone()[0]
    al_dia=c.execute("SELECT COUNT(*) FROM beneficiarios WHERE activo=1 AND cuotas_pagadas >0 AND cuotas_pagadas <11").fetchone()[0]
    pagados=c.execute("SELECT COUNT(*) FROM beneficiarios WHERE activo=1 AND cuotas_pagadas >=11").fetchone()[0]
    sectores=c.execute("SELECT sector,COUNT(*) cnt FROM beneficiarios WHERE activo=1 GROUP BY sector ORDER BY cnt DESC").fetchall()
    ult=c.execute("SELECT p.fecha_pago,p.monto_usd,b.nombres||' '||b.apellidos nombre FROM pagos p JOIN beneficiarios b ON p.beneficiario_id=b.id ORDER BY p.fecha_registro DESC LIMIT 5").fetchall()
    # Datos para alertas de RIF
    rif_inf = c.execute("SELECT id, nombres||' '||apellidos as nombre, contrato, rif_vencimiento FROM beneficiarios WHERE activo=1 AND rif_vencimiento IS NOT NULL").fetchall()
    
    conn.close()
    
    pct_rec = round((cobrado / cartera * 100), 2) if cartera > 0 else 0
    
    return jsonify({
        "total_beneficiarios":total, "cartera_total":cartera, "monto_cobrado":round(cobrado,2),
        "saldo_pendiente":round(cartera-cobrado,2), "en_mora":mora, "al_dia":al_dia, "pagados_completo":pagados,
        "porcentaje_recuperacion": pct_rec,
        "sectores":[{"sector":r["sector"] or "Sin sector", "count":r["cnt"]} for r in sectores],
        "ultimos_pagos":[dict(r) for r in ult],
        "rif_info": [dict(r) for r in rif_inf],
        "sync": "v2.3"
    })

# --- BENEFICIARIOS Y MAPA (CON ELIMINAR Y CORRECCIÓN DE MARCADOR) ---

@app.route('/api/beneficiarios')
@login_required
def api_beneficiarios():
    conn=get_db()
    rows=conn.execute("SELECT * FROM beneficiarios WHERE activo=1 ORDER BY contrato").fetchall()
    conn.close()
    return jsonify([dict(r) for r in rows])

@app.route('/api/beneficiarios/<int:bid>')
@login_required
def api_beneficiario(bid):
    conn=get_db()
    b=conn.execute("SELECT * FROM beneficiarios WHERE id=?",(bid,)).fetchone()
    pagos=conn.execute("SELECT * FROM pagos WHERE beneficiario_id=? ORDER BY numero_cuota",(bid,)).fetchall()
    conn.close()
    if not b: return jsonify({"error": "No encontrado"}),404
    return jsonify({"beneficiario":dict(b), "pagos":[dict(p) for p in pagos]})

# 🆕 NUEVA FUNCIÓN: ELIMINAR BENEFICIARIO (Soft Delete + Force)
@app.route('/api/beneficiarios/<int:bid>/eliminar', methods=['POST'])
@login_required
def eliminar_beneficiario(bid):
    d = request.json or {}
    force = d.get('force', False)

    conn = get_db()
    c = conn.cursor()
    
    # Verificar si tiene pagos registrados (nunca se permite, ni con force)
    pagos = c.execute("SELECT COUNT(*) FROM pagos WHERE beneficiario_id=?", (bid,)).fetchone()[0]
    if pagos > 0:
        conn.close()
        return jsonify({"error": "No se puede eliminar: Este beneficiario tiene pagos registrados. Imposible eliminar registros con historial financiero."}), 400
    
    # Verificar contratos generados
    contratos = c.execute("SELECT COUNT(*) FROM contratos_generados WHERE beneficiario_id=?", (bid,)).fetchone()[0]
    if contratos > 0 and not force:
        conn.close()
        return jsonify({
            "error_tipo": "contratos",
            "error": f"Este beneficiario tiene {contratos} contrato(s) registrado(s).",
            "contratos": contratos
        }), 400

    # Si force=True, eliminar primero los registros de contratos
    if force and contratos > 0:
        c.execute("DELETE FROM contratos_generados WHERE beneficiario_id=?", (bid,))

    # Soft Delete: Cambiar activo a 0
    c.execute("UPDATE beneficiarios SET activo=0 WHERE id=?", (bid,))
    conn.commit()
    conn.close()
    return jsonify({"ok": True, "mensaje": "Beneficiario eliminado correctamente."})

# 🆕 GEOCODIFICAR BENEFICIARIO desde su dirección (Nominatim/OpenStreetMap)
@app.route('/api/beneficiarios/<int:bid>/geocodificar', methods=['POST'])
@login_required
def geocodificar_beneficiario(bid):
    import urllib.request, urllib.parse, time
    conn = get_db()
    b = conn.execute("SELECT * FROM beneficiarios WHERE id=?", (bid,)).fetchone()
    conn.close()
    if not b:
        return jsonify({"error": "Beneficiario no encontrado"}), 404
    b = dict(b)

    cfg_conn = get_db()
    ciudad = cfg_conn.execute("SELECT valor FROM config_sistema WHERE clave='ciudad'").fetchone()
    estado = cfg_conn.execute("SELECT valor FROM config_sistema WHERE clave='estado'").fetchone()
    cfg_conn.close()
    ciudad = ciudad['valor'] if ciudad else ''
    estado = estado['valor'] if estado else 'Venezuela'

    direccion = (b.get('direccion_real') or b.get('direccion') or '').strip()
    sector    = (b.get('sector') or '').strip()

    partes = [p for p in [direccion, sector, ciudad, 'Venezuela'] if p]
    query  = ', '.join(partes)

    url = 'https://nominatim.openstreetmap.org/search?' + urllib.parse.urlencode(
        {'q': query, 'format': 'json', 'limit': 1, 'countrycodes': 've'}
    )
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'SemilleroComunal/1.2'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            results = json.loads(resp.read())
        if results:
            lat = float(results[0]['lat'])
            lng = float(results[0]['lon'])
            upd = get_db()
            upd.execute("UPDATE beneficiarios SET lat=?, lng=? WHERE id=?", (lat, lng, bid))
            upd.commit(); upd.close()
            return jsonify({"ok": True, "lat": lat, "lng": lng,
                            "address": results[0].get('display_name', ''), "query": query})
        else:
            return jsonify({"error": f"No se encontró ubicación para: {query}"}), 404
    except Exception as e:
        return jsonify({"error": f"Error de geocodificación: {str(e)}"}), 500

# 🆕 GEOCODIFICAR TODOS los beneficiarios sin coordenadas
@app.route('/api/beneficiarios/geocodificar_todos', methods=['POST'])
@login_required
def geocodificar_todos():
    import urllib.request, urllib.parse, time
    conn = get_db()
    # Beneficiarios sin coordenadas válidas (lat/lng == 0 o NULL o coordenada por defecto)
    bens = conn.execute(
        "SELECT * FROM beneficiarios WHERE activo=1 AND (lat IS NULL OR lng IS NULL OR (lat=8.623 AND lng=-71.651))"
    ).fetchall()
    cfg_ciudad = conn.execute("SELECT valor FROM config_sistema WHERE clave='ciudad'").fetchone()
    conn.close()
    ciudad = cfg_ciudad['valor'] if cfg_ciudad else ''

    resultados = []
    for b in bens:
        b = dict(b)
        direccion = (b.get('direccion_real') or b.get('direccion') or '').strip()
        sector    = (b.get('sector') or '').strip()
        partes = [p for p in [direccion, sector, ciudad, 'Venezuela'] if p]
        query  = ', '.join(partes)

        url = 'https://nominatim.openstreetmap.org/search?' + urllib.parse.urlencode(
            {'q': query, 'format': 'json', 'limit': 1, 'countrycodes': 've'}
        )
        try:
            req = urllib.request.Request(url, headers={'User-Agent': 'SemilleroComunal/1.2'})
            with urllib.request.urlopen(req, timeout=10) as resp:
                results = json.loads(resp.read())
            if results:
                lat = float(results[0]['lat'])
                lng = float(results[0]['lon'])
                upd = get_db()
                upd.execute("UPDATE beneficiarios SET lat=?, lng=? WHERE id=?", (lat, lng, b['id']))
                upd.commit(); upd.close()
                resultados.append({"id": b['id'], "nombre": b['nombres']+' '+b['apellidos'],
                                   "ok": True, "lat": lat, "lng": lng})
            else:
                resultados.append({"id": b['id'], "nombre": b['nombres']+' '+b['apellidos'],
                                   "ok": False, "error": "Dirección no encontrada"})
        except Exception as e:
            resultados.append({"id": b['id'], "nombre": b['nombres']+' '+b['apellidos'],
                               "ok": False, "error": str(e)})
        time.sleep(1.1)  # Respetar límite de Nominatim: 1 req/seg

    encontrados = sum(1 for r in resultados if r['ok'])
    return jsonify({"ok": True, "total": len(resultados), "encontrados": encontrados,
                   "resultados": resultados})

@app.route('/api/beneficiarios/<int:bid>/actualizar', methods=['POST'])
@login_required
def actualizar_beneficiario(bid):
    d=request.json
    conn=get_db()
    fields=["direccion_real", "sector", "actividad", "lat", "lng", "historia"]
    conn.execute(f"UPDATE beneficiarios SET {','.join(f+'=?' for f in fields)} WHERE id=?",
    [d.get(f) for f in fields]+[bid])
    conn.commit()
    conn.close()
    return jsonify({"ok":True})

@app.route('/api/beneficiarios/<int:bid>/ubicacion', methods=['POST'])
@login_required
def actualizar_ubicacion(bid):
    d=request.json
    conn=get_db()
    conn.execute("UPDATE beneficiarios SET lat=?,lng=?,direccion_real=? WHERE id=?",
    (d['lat'],d['lng'],d.get('direccion_real'),bid))
    conn.commit()
    conn.close()
    return jsonify({"ok":True})

@app.route('/api/beneficiarios/nuevo', methods=['POST'])
@login_required
def nuevo_beneficiario():
    d = request.json
    try:
        conn = get_db()
        c = conn.cursor()
        
        # ELIMINADO DUPLICADOS POTENCIALES: Si ya existe un beneficiario con esta cédula activo, abortar
        existe = c.execute("SELECT id FROM beneficiarios WHERE cedula=? AND activo=1", (d.get('cedula'),)).fetchone()
        if existe:
            conn.close()
            return jsonify({"error": f"Ya existe un beneficiario activo con la cédula {d.get('cedula')}"}), 400

        # LÓGICA DE NUMERACIÓN PROFUNDA:
        # 1. Buscamos el número más alto entre los beneficiarios ACTIVOS
        activos = c.execute("SELECT contrato FROM beneficiarios WHERE activo=1").fetchall()
        max_num = 0
        for row in activos:
            try:
                partes = row['contrato'].split('-')
                if len(partes) >= 2:
                    n = int(partes[1])
                    if n > max_num: max_num = n
            except: pass
        
        num = max_num + 1
        año = date.today().year
        contrato = f"BCCPPO-{num:02d}-{año}"
        
        # 2. LIMPIEZA: Si existe un registro INACTIVO con el mismo número de contrato, lo eliminamos 
        # (ya que el usuario quiere reutilizar números de registros que borró por error)
        c.execute("DELETE FROM beneficiarios WHERE contrato=? AND activo=0", (contrato,))
        
        mun = get_cfg('municipio') or ''
        coords = MUNICIPIOS_COORDS.get(mun, [8.623, -71.651])
        
        lat = d.get('lat')
        lng = d.get('lng')
        if not lat or str(lat).strip() == "" or not lng or str(lng).strip() == "":
            lat, lng = coords[0], coords[1]
        
        c.execute("""INSERT INTO beneficiarios(contrato,nombres,apellidos,cedula,rif,direccion,sector,lat,lng,
        estado_civil,fecha_nacimiento,actividad,fecha_desembolso,rif_vencimiento)
        VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
        (contrato, d.get('nombres','').strip(), d.get('apellidos','').strip(), d.get('cedula','').strip(), 
         d.get('rif','').strip(), d.get('direccion','').strip(), d.get('sector','').strip(),
         lat, lng, d.get('estado_civil',''), d.get('fecha_nacimiento',''),
         d.get('actividad',''), d.get('fecha_desembolso', date.today().strftime('%d/%m/%Y')), d.get('rif_vencimiento','')))
        
        nid = c.lastrowid
        c.execute("INSERT INTO movimientos_banco(tipo,concepto,monto_usd,fecha,referencia,beneficiario_id) VALUES('EGRESO',?,400,?,?,?)",
        (f"Desembolso {contrato}", d.get('fecha_desembolso', date.today().strftime('%d/%m/%Y')), contrato, nid))
        
        conn.commit()
        conn.close()
        logger.info(f"Beneficiario creado: {contrato} ID:{nid}")
        return jsonify({"ok": True, "contrato": contrato, "id": nid})
    except Exception as e:
        logger.error(f"Error crítico en nuevo_beneficiario: {e}")
        return jsonify({"error": str(e)}), 500

# --- PAGOS Y BANCO ---

@app.route('/api/pagos/registrar', methods=['POST'])
@login_required
def registrar_pago():
    d=request.json
    conn=get_db()
    c=conn.cursor()
    b=c.execute("SELECT * FROM beneficiarios WHERE id=?",(d['beneficiario_id'],)).fetchone()
    if not b: return jsonify({"error": "No encontrado"}),404
    nueva=b['cuotas_pagadas']+1
    if nueva > b['cuotas_total']: return jsonify({"error": "Ya completó todas las cuotas"}),400
    
    c.execute("""INSERT INTO pagos(beneficiario_id,contrato,numero_cuota,monto_usd,tasa_bcv,monto_bs,
    fecha_pago,numero_operacion,banco_origen,foto_capture,registrado_por) VALUES(?,?,?,?,?,?,?,?,?,?,?)""",
    (d['beneficiario_id'],b['contrato'],nueva,d.get('monto_usd',40),d.get('tasa_bcv'),d.get('monto_bs'),
    d.get('fecha_pago'),d.get('numero_operacion'),d.get('banco_origen'),d.get('foto_capture'),session.get('usuario','Sistema')))
    
    c.execute("UPDATE beneficiarios SET cuotas_pagadas=cuotas_pagadas+1 WHERE id=?",(d['beneficiario_id'],))
    c.execute("INSERT INTO movimientos_banco(tipo,concepto,monto_usd,fecha,referencia,beneficiario_id) VALUES('INGRESO',?,?,?,?,?)",
    (f"Cuota {nueva}/11",d.get('monto_usd',40),d.get('fecha_pago'),d.get('numero_operacion'),d['beneficiario_id']))
    conn.commit()
    conn.close()
    return jsonify({"ok":True, "cuota":nueva})

@app.route('/api/banco/movimiento', methods=['POST'])
@login_required
def agregar_movimiento():
    d=request.json
    conn=get_db()
    conn.execute("INSERT INTO movimientos_banco(tipo,concepto,monto_usd,fecha,referencia,notas) VALUES(?,?,?,?,?,?)",
    (d['tipo'],d['concepto'],float(d['monto_usd']),d['fecha'],d.get('referencia',''),d.get('notas','')))
    conn.commit()
    conn.close()
    return jsonify({"ok":True})

@app.route('/api/estado_banco')
@login_required
def estado_banco():
    conn=get_db()
    ing=conn.execute("SELECT SUM(monto_usd) FROM movimientos_banco WHERE tipo='INGRESO'").fetchone()[0] or 0
    egr=conn.execute("SELECT SUM(monto_usd) FROM movimientos_banco WHERE tipo='EGRESO'").fetchone()[0] or 0
    movs=conn.execute("SELECT * FROM movimientos_banco ORDER BY fecha_registro DESC LIMIT 30").fetchall()
    conn.close()
    return jsonify({"total_ingresos":round(ing,2),"total_egresos":round(egr,2),"saldo":round(ing-egr,2),"movimientos":[dict(m) for m in movs]})

# --- FOTOS E HISTORIA ---

@app.route('/api/fotos/agregar', methods=['POST'])
@login_required
def agregar_foto():
    d=request.json
    conn=get_db()
    conn.execute("INSERT INTO fotos_historia(beneficiario_id,imagen_b64,descripcion,fecha,tipo) VALUES(?,?,?,?,?)",
    (d['beneficiario_id'],d['imagen_b64'],d.get('descripcion',''),d.get('fecha',date.today().strftime('%d/%m/%Y')),d.get('tipo','historia')))
    conn.commit()
    fid=conn.execute("SELECT last_insert_rowid()").fetchone()[0]
    conn.close()
    return jsonify({"ok":True,"id":fid})

@app.route('/api/fotos/<int:fid>/img')
@login_required
def get_foto_img(fid):
    conn=get_db()
    f=conn.execute("SELECT imagen_b64 FROM fotos_historia WHERE id=?",(fid,)).fetchone()
    conn.close()
    if not f: return '',404
    try: return Response(base64.b64decode(f['imagen_b64']),mimetype='image/jpeg')
    except: return '',500

@app.route('/api/fotos/<int:fid>/eliminar', methods=['POST'])
@login_required
def eliminar_foto(fid):
    conn=get_db()
    conn.execute("DELETE FROM fotos_historia WHERE id=?",(fid,))
    conn.commit()
    conn.close()
    return jsonify({"ok":True})

# --- OCR MULTI-PROVEEDOR Y CONCILIACIÓN BANCARIA ---

def get_gemini():
    import google.generativeai as genai
    ak=get_cfg('api_key')
    if not ak: return None, "API Key no configurada"
    genai.configure(api_key=ak)
    return genai.GenerativeModel('gemini-2.5-flash'), None

def call_groq_vision(image_b64, prompt):
    import urllib.request
    import json
    
    g_key = get_cfg('groq_api_key')
    if not g_key:
        raise ValueError("Groq API Key no configurada")
        
    url = "https://api.groq.com/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {g_key}",
        "Content-Type": "application/json"
    }
    
    data = {
        "model": "llama-3.2-11b-vision-preview",
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": prompt
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{image_b64}"
                        }
                    }
                ]
            }
        ],
        "temperature": 0.1,
        "response_format": {
            "type": "json_object"
        }
    }
    
    req = urllib.request.Request(url, data=json.dumps(data).encode('utf-8'), headers=headers, method='POST')
    
    try:
        with urllib.request.urlopen(req, timeout=30) as response:
            res_body = response.read().decode('utf-8')
            res_json = json.loads(res_body)
            content = res_json['choices'][0]['message']['content']
            return content
    except Exception as e:
        logger.error(f"Groq Vision API error: {e}")
        raise

def process_ocr_pago(image_b64):
    pref = get_cfg('api_preferida') or 'gemini'
    prompt = 'Analiza este comprobante de pago bancario venezolano. Extrae la informacion y responde estrictamente en un objeto JSON sin formato adicional con las siguientes claves: {"monto_bs": "", "monto_usd": "", "fecha": "DD/MM/AAAA", "numero_operacion": "", "banco_origen": "", "tasa_bcv": ""}. Trata de limpiar caracteres no numericos en el monto (ej. extraer 20800.00 en vez de Bs. 20.800).'
    
    errors = []
    if pref == 'groq':
        try:
            res_text = call_groq_vision(image_b64, prompt)
            return json.loads(res_text)
        except Exception as e:
            errors.append(f"Groq error: {e}")
            logger.warning(f"OCR Pago: Groq falló, intentando Gemini. Error: {e}")
            try:
                model, err = get_gemini()
                if not err:
                    r = model.generate_content([
                        {"inline_data": {"mime_type": "image/jpeg", "data": image_b64}},
                        {"text": prompt}
                    ])
                    text = r.text.replace('```json','').replace('```','').strip()
                    return json.loads(text)
                else:
                    errors.append(f"Gemini init error: {err}")
            except Exception as e2:
                errors.append(f"Gemini failover error: {e2}")
    else:
        try:
            model, err = get_gemini()
            if err: raise ValueError(err)
            r = model.generate_content([
                {"inline_data": {"mime_type": "image/jpeg", "data": image_b64}},
                {"text": prompt}
            ])
            text = r.text.replace('```json','').replace('```','').strip()
            return json.loads(text)
        except Exception as e:
            errors.append(f"Gemini error: {e}")
            logger.warning(f"OCR Pago: Gemini falló, intentando Groq. Error: {e}")
            try:
                res_text = call_groq_vision(image_b64, prompt)
                return json.loads(res_text)
            except Exception as e2:
                errors.append(f"Groq failover error: {e2}")
                
    raise ValueError(f"Ambos proveedores de OCR fallaron. Detalles: {', '.join(errors)}")

def process_ocr_docs(cedula_b64, rif_b64):
    pref = get_cfg('api_preferida') or 'gemini'
    prompt = 'Analiza estos documentos de identidad venezolanos (Cédula de Identidad y/o RIF). Extrae los datos y responde estrictamente en un objeto JSON sin formato adicional con las siguientes claves: {"nombres": "", "apellidos": "", "cedula": "V-X.XXX.XXX", "rif": "VXXXXXXXXX", "estado_civil": "", "fecha_nacimiento": "DD/MM/AAAA", "rif_vencimiento": "DD/MM/AAAA", "direccion": ""}'
    
    errors = []
    img_b64 = cedula_b64 or rif_b64
    
    if pref == 'groq':
        try:
            res_text = call_groq_vision(img_b64, prompt)
            return json.loads(res_text)
        except Exception as e:
            errors.append(f"Groq error: {e}")
            logger.warning(f"OCR Docs: Groq falló, intentando Gemini. Error: {e}")
            try:
                model, err = get_gemini()
                if not err:
                    parts = []
                    if cedula_b64: parts.append({"inline_data": {"mime_type": "image/jpeg", "data": cedula_b64}})
                    if rif_b64: parts.append({"inline_data": {"mime_type": "image/jpeg", "data": rif_b64}})
                    parts.append({"text": prompt})
                    r = model.generate_content(parts)
                    text = r.text.replace('```json','').replace('```','').strip()
                    return json.loads(text)
                else:
                    errors.append(f"Gemini init error: {err}")
            except Exception as e2:
                errors.append(f"Gemini failover error: {e2}")
    else:
        try:
            model, err = get_gemini()
            if err: raise ValueError(err)
            parts = []
            if cedula_b64: parts.append({"inline_data": {"mime_type": "image/jpeg", "data": cedula_b64}})
            if rif_b64: parts.append({"inline_data": {"mime_type": "image/jpeg", "data": rif_b64}})
            parts.append({"text": prompt})
            r = model.generate_content(parts)
            text = r.text.replace('```json','').replace('```','').strip()
            return json.loads(text)
        except Exception as e:
            errors.append(f"Gemini error: {e}")
            logger.warning(f"OCR Docs: Gemini falló, intentando Groq. Error: {e}")
            try:
                res_text = call_groq_vision(img_b64, prompt)
                return json.loads(res_text)
            except Exception as e2:
                errors.append(f"Groq failover error: {e2}")
                
    raise ValueError(f"Ambos proveedores de OCR fallaron. Detalles: {', '.join(errors)}")

def process_ocr_estado(image_b64):
    pref = get_cfg('api_preferida') or 'gemini'
    prompt = 'Analiza este capture de estado de cuenta bancario venezolano. Extrae la lista de transacciones y responde estrictamente en un objeto JSON con una única clave "transacciones" que contenga una lista de objetos. Cada objeto debe tener los siguientes campos: {"fecha": "DD/MM/AAAA", "referencia": "", "descripcion": "", "codigo": "NC o ND", "debito": 0.0, "credito": 0.0, "saldo": 0.0}. Limpia los montos para que sean números flotantes puros (sin símbolos de moneda ni separadores de miles).'
    
    errors = []
    if pref == 'groq':
        try:
            res_text = call_groq_vision(image_b64, prompt)
            return json.loads(res_text)
        except Exception as e:
            errors.append(f"Groq error: {e}")
            logger.warning(f"OCR Estado: Groq falló, intentando Gemini. Error: {e}")
            try:
                model, err = get_gemini()
                if not err:
                    r = model.generate_content([
                        {"inline_data": {"mime_type": "image/jpeg", "data": image_b64}},
                        {"text": prompt}
                    ])
                    text = r.text.replace('```json','').replace('```','').strip()
                    return json.loads(text)
            except Exception as e2:
                errors.append(f"Gemini failover error: {e2}")
    else:
        try:
            model, err = get_gemini()
            if err: raise ValueError(err)
            r = model.generate_content([
                {"inline_data": {"mime_type": "image/jpeg", "data": image_b64}},
                {"text": prompt}
            ])
            text = r.text.replace('```json','').replace('```','').strip()
            return json.loads(text)
        except Exception as e:
            errors.append(f"Gemini error: {e}")
            logger.warning(f"OCR Estado: Gemini falló, intentando Groq. Error: {e}")
            try:
                res_text = call_groq_vision(image_b64, prompt)
                return json.loads(res_text)
            except Exception as e2:
                errors.append(f"Groq failover error: {e2}")
                
    raise ValueError(f"Ambos proveedores de OCR fallaron para el estado de cuenta. Detalles: {', '.join(errors)}")

def parse_bdt_statement_text(text):
    import re
    lines = text.split('\n')
    transacciones = []
    
    for line in lines:
        line = line.strip()
        if not line: continue
        if "referencia" in line.lower() or "saldo" in line.lower() or "estado de cuenta" in line.lower():
            continue
            
        fecha_match = re.match(r'^(\d{2}[/-]\d{2}[/-]\d{4})', line)
        if not fecha_match:
            continue
            
        fecha = fecha_match.group(1)
        rest = line[len(fecha):].strip()
        
        tokens = [t.strip() for t in re.split(r'\s{2,}', rest) if t.strip()]
        if len(tokens) < 3:
            tokens = [t.strip() for t in rest.split(' ') if t.strip()]
            
        if len(tokens) >= 3:
            referencia = tokens[0]
            codigo = tokens[1] if len(tokens[1]) <= 3 else ""
            
            def parse_monto(val):
                if not val or val in ('0', '0,00', '0.00', '-'): 
                    return 0.0
                val_clean = val.replace('.', '').replace(',', '.')
                try:
                    return float(val_clean)
                except:
                    val_clean = re.sub(r'[^\d.]', '', val.replace(',', '.'))
                    try: return float(val_clean)
                    except: return 0.0
            
            saldo = parse_monto(tokens[-1])
            credito = parse_monto(tokens[-2])
            debito = parse_monto(tokens[-3])
            
            desc_start = 2 if codigo else 1
            desc_tokens = tokens[desc_start:-3]
            descripcion = " ".join(desc_tokens) if desc_tokens else "TRANSFERENCIA"
            
            transacciones.append({
                "fecha": fecha,
                "referencia": referencia,
                "descripcion": descripcion,
                "codigo": codigo,
                "debito": debito,
                "credito": credito,
                "saldo": saldo
            })
            
    return transacciones

@app.route('/api/ocr', methods=['POST'])
@login_required
def ocr_docs():
    try:
        cedula_b64 = request.json.get('cedula_b64')
        rif_b64 = request.json.get('rif_b64')
        res = process_ocr_docs(cedula_b64, rif_b64)
        return jsonify(res)
    except Exception as e:
        logger.error(f"OCR docs error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/ocr/pago', methods=['POST'])
@login_required
def ocr_pago():
    try:
        img_b64 = request.json.get('imagen_b64')
        if not img_b64: return jsonify({"error": "No se proporcionó imagen"}), 400
        res = process_ocr_pago(img_b64)
        return jsonify(res)
    except Exception as e:
        logger.error(f"OCR pago error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/conciliacion/subir_estado', methods=['POST'])
@login_required
def subir_estado_cuenta():
    d = request.json or {}
    text_data = d.get('texto')
    image_b64 = d.get('imagen_b64')
    
    transacciones = []
    try:
        if text_data:
            transacciones = parse_bdt_statement_text(text_data)
        elif image_b64:
            res = process_ocr_estado(image_b64)
            transacciones = res.get('transacciones', [])
        else:
            return jsonify({"error": "No se proporcionó texto ni imagen del estado de cuenta"}), 400
            
        conn = get_db()
        c = conn.cursor()
        
        insertadas = 0
        duplicadas = 0
        
        for tx in transacciones:
            fecha = tx.get('fecha')
            referencia = tx.get('referencia')
            descripcion = tx.get('descripcion', '')
            codigo = tx.get('codigo', '')
            debito = tx.get('debito', 0.0)
            credito = tx.get('credito', 0.0)
            saldo = tx.get('saldo', 0.0)
            
            if not referencia or not fecha:
                continue
                
            try:
                c.execute("""INSERT INTO transacciones_banco
                (fecha, referencia, descripcion, codigo, debito, credito, saldo, reconciliado)
                VALUES(?,?,?,?,?,?,?,0)""",
                (fecha, referencia, descripcion, codigo, debito, credito, saldo))
                insertadas += 1
            except sqlite3.IntegrityError:
                duplicadas += 1
                
        conn.commit()
        conn.close()
        
        return jsonify({
            "ok": True,
            "mensaje": f"Se procesaron {len(transacciones)} transacciones. Insertadas: {insertadas}, Duplicadas (omitidas): {duplicadas}",
            "transacciones_procesadas": len(transacciones),
            "insertadas": insertadas,
            "duplicadas": duplicadas
        })
    except Exception as e:
        logger.error(f"Error al subir estado de cuenta: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/conciliacion/ejecutar', methods=['POST'])
@login_required
def ejecutar_conciliacion():
    try:
        conn = get_db()
        c = conn.cursor()
        
        pagos = conn.execute("SELECT id, contrato, numero_cuota, monto_bs, numero_operacion, fecha_pago FROM pagos WHERE conciliado=0").fetchall()
        bancos = conn.execute("SELECT id, fecha, referencia, descripcion, credito FROM transacciones_banco WHERE reconciliado=0 AND credito > 0").fetchall()
        
        emparejados = 0
        for p in pagos:
            p_id = p['id']
            p_ref = str(p['numero_operacion'] or '').strip()
            p_monto = p['monto_bs'] or 0.0
            
            if not p_ref or p_monto == 0:
                continue
                
            match_encontrado = None
            for b in bancos:
                b_id = b['id']
                b_ref = str(b['referencia'] or '').strip()
                b_monto = b['credito'] or 0.0
                
                # Tolerancia de variación de monto de 100 Bs
                if abs(p_monto - b_monto) <= 100.0:
                    match_refs = False
                    if p_ref == b_ref:
                        match_refs = True
                    elif len(p_ref) >= 6 and len(b_ref) >= 6:
                        if p_ref.endswith(b_ref) or b_ref.endswith(p_ref):
                            match_refs = True
                            
                    if match_refs:
                        match_encontrado = b
                        break
            
            if match_encontrado:
                b_id = match_encontrado['id']
                b_ref = match_encontrado['referencia']
                
                c.execute("UPDATE pagos SET conciliado=1, banco_id_referencia=? WHERE id=?", (b_ref, p_id))
                c.execute("UPDATE transacciones_banco SET reconciliado=1, pago_id=? WHERE id=?", (p_id, b_id))
                emparejados += 1
                bancos = [x for x in bancos if x['id'] != b_id]
                
        conn.commit()
        conn.close()
        
        return jsonify({
            "ok": True,
            "mensaje": f"Conciliación finalizada. Se auto-conciliaron {emparejados} pagos con sus transacciones bancarias.",
            "emparejados": emparejados
        })
    except Exception as e:
        logger.error(f"Error en ejecución de conciliación: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/conciliacion/manual', methods=['POST'])
@login_required
def conciliacion_manual():
    d = request.json or {}
    pago_id = d.get('pago_id')
    transaccion_id = d.get('transaccion_id')
    
    if not pago_id or not transaccion_id:
        return jsonify({"error": "Faltan IDs de pago o transacción"}), 400
        
    try:
        conn = get_db()
        c = conn.cursor()
        
        pago = conn.execute("SELECT * FROM pagos WHERE id=?", (pago_id,)).fetchone()
        if not pago:
            conn.close()
            return jsonify({"error": "Pago no encontrado"}), 404
            
        tx = conn.execute("SELECT * FROM transacciones_banco WHERE id=?", (transaccion_id,)).fetchone()
        if not tx:
            conn.close()
            return jsonify({"error": "Transacción bancaria no encontrada"}), 404
            
        if pago['banco_id_referencia']:
            c.execute("UPDATE transacciones_banco SET reconciliado=0, pago_id=NULL WHERE referencia=?", (pago['banco_id_referencia'],))
        if tx['pago_id']:
            c.execute("UPDATE pagos SET conciliado=0, banco_id_referencia=NULL WHERE id=?", (tx['pago_id'],))
            
        c.execute("UPDATE pagos SET conciliado=2, banco_id_referencia=? WHERE id=?", (tx['referencia'], pago_id))
        c.execute("UPDATE transacciones_banco SET reconciliado=2, pago_id=? WHERE id=?", (pago_id, transaccion_id))
        
        conn.commit()
        conn.close()
        return jsonify({"ok": True, "mensaje": "Conciliación manual realizada exitosamente."})
    except Exception as e:
        logger.error(f"Error en conciliación manual: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/conciliacion/estado', methods=['GET'])
@login_required
def estado_conciliacion():
    try:
        conn = get_db()
        txs = conn.execute("SELECT * FROM transacciones_banco ORDER BY fecha DESC, id DESC").fetchall()
        pagos = conn.execute("""
            SELECT p.*, b.nombres, b.apellidos 
            FROM pagos p 
            JOIN beneficiarios b ON p.beneficiario_id = b.id 
            ORDER BY p.fecha_pago DESC, p.id DESC
        """).fetchall()
        conn.close()
        return jsonify({
            "transacciones": [dict(x) for x in txs],
            "pagos": [dict(x) for x in pagos]
        })
    except Exception as e:
        logger.error(f"Error al obtener estado de conciliación: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/pagos/cuadro_resumen', methods=['GET'])
@login_required
def cuadro_resumen():
    try:
        conn = get_db()
        beneficiarios = conn.execute("SELECT id, contrato, nombres, apellidos, cedula FROM beneficiarios WHERE activo=1 ORDER BY contrato").fetchall()
        pagos = conn.execute("SELECT id, beneficiario_id, numero_cuota, monto_bs, monto_usd, fecha_pago, numero_operacion, banco_origen, conciliado FROM pagos").fetchall()
        conn.close()
        
        pagos_map = {}
        for p in pagos:
            b_id = p['beneficiario_id']
            cuota = p['numero_cuota']
            if b_id not in pagos_map:
                pagos_map[b_id] = {}
            pagos_map[b_id][cuota] = dict(p)
            
        matriz = []
        for b in beneficiarios:
            b_id = b['id']
            cuotas_info = {}
            for c_num in range(1, 12):
                pago_cuota = pagos_map.get(b_id, {}).get(c_num)
                if pago_cuota:
                    status = "pendiente"
                    if pago_cuota['conciliado'] == 1:
                        status = "conciliado_auto"
                    elif pago_cuota['conciliado'] == 2:
                        status = "conciliado_manual"
                    else:
                        status = "registrado"
                        
                    cuotas_info[str(c_num)] = {
                        "status": status,
                        "monto_bs": pago_cuota['monto_bs'],
                        "monto_usd": pago_cuota['monto_usd'],
                        "fecha": pago_cuota['fecha_pago'],
                        "referencia": pago_cuota['numero_operacion'],
                        "banco": pago_cuota['banco_origen']
                    }
                else:
                    cuotas_info[str(c_num)] = {
                        "status": "pendiente",
                        "monto_bs": 0,
                        "monto_usd": 0,
                        "fecha": "",
                        "referencia": "",
                        "banco": ""
                    }
            matriz.append({
                "beneficiario_id": b_id,
                "contrato": b['contrato'],
                "nombres": b['nombres'],
                "apellidos": b['apellidos'],
                "cedula": b['cedula'],
                "cuotas": cuotas_info
            })
        return jsonify(matriz)
    except Exception as e:
        logger.error(f"Error en cuadro_resumen: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/pagos/informe_pdf', methods=['GET'])
@login_required
def generar_reporte_pdf():
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        
        conn = get_db()
        beneficiarios = conn.execute("SELECT id, contrato, nombres, apellidos, cedula FROM beneficiarios WHERE activo=1 ORDER BY contrato").fetchall()
        pagos = conn.execute("SELECT id, beneficiario_id, numero_cuota, monto_bs, monto_usd, fecha_pago, numero_operacion, banco_origen, conciliado FROM pagos").fetchall()
        txs = conn.execute("SELECT id, fecha, referencia, descripcion, credito, reconciliado FROM transacciones_banco WHERE credito > 0").fetchall()
        conn.close()
        
        total_beneficiarios = len(beneficiarios)
        total_pagos_reg = len(pagos)
        pagos_conciliados = sum(1 for p in pagos if p['conciliado'] > 0)
        pagos_pendientes = total_pagos_reg - pagos_conciliados
        monto_total_bs = sum(p['monto_bs'] or 0.0 for p in pagos)
        monto_conciliado_bs = sum(p['monto_bs'] or 0.0 for p in pagos if p['conciliado'] > 0)
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
        story = []
        
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'ComunalTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=20,
            textColor=colors.HexColor('#0C447C'),
            spaceAfter=15,
            alignment=1
        )
        
        subtitle_style = ParagraphStyle(
            'ComunalSubTitle',
            parent=styles['Normal'],
            fontName='Helvetica-Oblique',
            fontSize=10,
            textColor=colors.HexColor('#555555'),
            spaceAfter=20,
            alignment=1
        )
        
        h2_style = ParagraphStyle(
            'ComunalH2',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=13,
            textColor=colors.HexColor('#0F6E56'),
            spaceBefore=12,
            spaceAfter=8
        )
        
        body_style = ParagraphStyle(
            'ComunalBody',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=9,
            textColor=colors.HexColor('#333333'),
            spaceAfter=8
        )
        
        story.append(Paragraph("BANCO COMUNAL MANUELA SÁENZ 5", title_style))
        story.append(Paragraph("Informe de Conciliación Bancaria y Registro de Pagos - Semillero Comunal", subtitle_style))
        
        story.append(Paragraph("Resumen Ejecutivo", h2_style))
        
        porcentaje_conciliado = (pagos_conciliados / total_pagos_reg * 100) if total_pagos_reg else 0.0
        data_resumen = [
            [Paragraph("<b>Indicador</b>", body_style), Paragraph("<b>Valor</b>", body_style)],
            [Paragraph("Beneficiarios Activos", body_style), Paragraph(str(total_beneficiarios), body_style)],
            [Paragraph("Total Pagos Registrados", body_style), Paragraph(str(total_pagos_reg), body_style)],
            [Paragraph("Pagos Conciliados con Banco", body_style), Paragraph(f"{pagos_conciliados} ({porcentaje_conciliado:.1f}%)", body_style)],
            [Paragraph("Pagos Pendientes de Conciliación", body_style), Paragraph(str(pagos_pendientes), body_style)],
            [Paragraph("Monto Total Recibido (Bs.)", body_style), Paragraph(f"{monto_total_bs:,.2f} Bs.", body_style)],
            [Paragraph("Monto Conciliado en Banco (Bs.)", body_style), Paragraph(f"{monto_conciliado_bs:,.2f} Bs.", body_style)]
        ]
        
        t_resumen = Table(data_resumen, colWidths=[250, 250])
        t_resumen.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#ECEFF1')),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CFD8DC')),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
            ('TOPPADDING', (0,0), (-1,-1), 5),
        ]))
        story.append(t_resumen)
        story.append(Spacer(1, 15))
        
        story.append(Paragraph("Cuadro de Pagos por Beneficiario", h2_style))
        
        pagos_map_count = {}
        for p in pagos:
            b_id = p['beneficiario_id']
            if b_id not in pagos_map_count:
                pagos_map_count[b_id] = 0
            if p['conciliado'] > 0:
                pagos_map_count[b_id] += 1
                
        data_beneficiarios = [
            [
                Paragraph("<b>Contrato</b>", body_style),
                Paragraph("<b>Cédula</b>", body_style),
                Paragraph("<b>Nombre Completo</b>", body_style),
                Paragraph("<b>Cuotas Conciliadas</b>", body_style),
                Paragraph("<b>Total Conciliado (Bs.)</b>", body_style)
            ]
        ]
        
        for b in beneficiarios:
            b_id = b['id']
            cuotas_pagadas_conc = pagos_map_count.get(b_id, 0)
            monto_beneficiario = sum(p['monto_bs'] or 0.0 for p in pagos if p['beneficiario_id'] == b_id and p['conciliado'] > 0)
            
            data_beneficiarios.append([
                Paragraph(b['contrato'], body_style),
                Paragraph(b['cedula'], body_style),
                Paragraph(f"{b['nombres']} {b['apellidos']}", body_style),
                Paragraph(f"{cuotas_pagadas_conc} / 11", body_style),
                Paragraph(f"{monto_beneficiario:,.2f} Bs.", body_style)
            ])
            
        t_benef = Table(data_beneficiarios, colWidths=[100, 80, 180, 100, 110])
        t_benef.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#0C447C')),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#B0BEC5')),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ('TOPPADDING', (0,0), (-1,-1), 4),
        ]))
        
        for col_idx in range(5):
            t_benef.setStyle(TableStyle([('TEXTCOLOR', (col_idx, 0), (col_idx, 0), colors.white)]))
            
        story.append(t_benef)
        story.append(Spacer(1, 15))
        
        story.append(Paragraph(f"Reporte generado el: {datetime.now().strftime('%d/%m/%Y %I:%M %p')}", subtitle_style))
        
        doc.build(story)
        buffer.seek(0)
        
        return send_file(
            buffer,
            as_attachment=True,
            download_name=f"Reporte_Conciliacion_{datetime.now().strftime('%Y%m%d')}.pdf",
            mimetype="application/pdf"
        )
    except Exception as e:
        logger.error(f"Error al generar reporte PDF: {e}")
        return jsonify({"error": str(e)}), 500

# --- CONTRATOS ---

@app.route('/api/contratos/generar', methods=['POST'])
@login_required
def generar_contrato():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    
    d=request.json
    bid=d.get('beneficiario_id')
    conn=get_db()
    b=conn.execute("SELECT * FROM beneficiarios WHERE id=?",(bid,)).fetchone()
    voceras=conn.execute("SELECT * FROM voceras WHERE activo=1 ORDER BY orden").fetchall()
    cfg={r['clave']:r['valor'] for r in conn.execute("SELECT clave,valor FROM config_sistema").fetchall()}
    conn.close()
    
    if not b: return jsonify({"error": "Beneficiario no encontrado"}),404

    nombre_banco=cfg.get('nombre_banco','BANCO DE LA COMUNA')
    ciudad=cfg.get('ciudad',cfg.get('municipio',''))
    estado=cfg.get('estado','')

    plantilla_path=os.path.join(BASE_DIR,'plantillas','modelo_contrato.docx')
    if os.path.exists(plantilla_path) and d.get('usar_modelo_propio'):
        doc=Document(plantilla_path)
        replacements={
            '[NOMBRE_BENEFICIARIO]': f"{b['nombres']} {b['apellidos']}",
            '[CEDULA_BENEFICIARIO]': b['cedula'],
            '[RIF_BENEFICIARIO]': b['rif'],
            '[NUMERO_CONTRATO]': b['contrato'],
            '[CIUDAD]': ciudad,
            '[ESTADO]': estado,
        }
        for para in doc.paragraphs:
            for k,v in replacements.items():
                if k in para.text:
                    for run in para.runs:
                        run.text=run.text.replace(k,v)
    else:
        doc=Document()
        for section in doc.sections:
            section.top_margin=Inches(1); section.bottom_margin=Inches(1)
            section.left_margin=Inches(1.2); section.right_margin=Inches(1.2)
        
        t=doc.add_paragraph()
        t.alignment=WD_ALIGN_PARAGRAPH.CENTER
        run=t.add_run('CONTRATO DE FINANCIAMIENTO PRODUCTIVO COMUNAL')
        run.bold=True; run.font.size=Pt(13)
        
        cn=doc.add_paragraph()
        cn.alignment=WD_ALIGN_PARAGRAPH.CENTER
        cnr=cn.add_run(b['contrato']); cnr.bold=True; cnr.font.size=Pt(12)
        
        doc.add_paragraph()
        voceras_text=' '.join([f"{v['nombres']} {v['apellidos']}, C.I. Nº {v['cedula']}." for v in voceras])
        
        body_text=(
            f"Entre {nombre_banco.upper()}, actuando mediante sus voceros: {voceras_text} (EL BANCO); "
            f"y {b['nombres']} {b['apellidos']}, C.I. {b['cedula']} RIF: {b['rif']} (EL BENEFICIARIO). "
            f"CLÁUSULA PRIMERA: OBJETO. Otorga financiamiento de USD $400. "
            f"CLÁUSULA SEGUNDA: PAGO. Once (11) cuotas mensuales de USD $40. "
            f"CLÁUSULA TERCERA: INCUMPLIMIENTO. Genera reporte negativo y suspensión de beneficios. "
            f"CLÁUSULA CUARTA: JURISDICCIÓN. Tribunales de {ciudad}, Estado {estado}."
        )
        bp=doc.add_paragraph()
        bp.alignment=WD_ALIGN_PARAGRAPH.JUSTIFY
        bp.add_run(body_text).font.size=Pt(10)
        
        doc.add_paragraph()
        sp=doc.add_paragraph(); sp.add_run('POR EL BANCO:').bold=True
        for v in voceras:
            doc.add_paragraph(f"{v['nombres']} {v['apellidos']} - C.I. {v['cedula']}")
        
        doc.add_paragraph()
        bfp=doc.add_paragraph(); bfp.add_run('EL BENEFICIARIO:').bold=True
        doc.add_paragraph(f"{b['nombres']} {b['apellidos']} - C.I. {b['cedula']}")

    fname=f"Contrato_{b['apellidos'].replace(' ','_')}_{b['contrato']}.docx"
    fpath=os.path.join(BASE_DIR,'documentos',fname)
    doc.save(fpath)
    
    conn=get_db()
    conn.execute("INSERT INTO contratos_generados(beneficiario_id,contrato,ruta_archivo) VALUES(?,?,?)",(bid,b['contrato'],fname))
    conn.commit()
    conn.close()
    
    return send_file(fpath, as_attachment=True, download_name=fname, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

@app.route('/api/contratos/subir_modelo', methods=['POST'])
@login_required
def subir_modelo():
    if 'archivo' not in request.files: return jsonify({"error":"Sin archivo"}),400
    f=request.files['archivo']
    if not f.filename.endswith('.docx'): return jsonify({"error":"Solo archivos .docx"}),400
    fpath=os.path.join(BASE_DIR,'plantillas','modelo_contrato.docx')
    f.save(fpath)
    return jsonify({"ok":True,"mensaje":"Modelo guardado"})

# --- REPORTES PDF Y PPTX ---

@app.route('/api/informe/pdf')
@login_required
def generar_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    
    conn=get_db()
    cfg={r['clave']:r['valor'] for r in conn.execute("SELECT clave,valor FROM config_sistema").fetchall()}
    total=conn.execute("SELECT COUNT(*) FROM beneficiarios WHERE activo=1").fetchone()[0]
    cartera=conn.execute("SELECT SUM(monto) FROM beneficiarios WHERE activo=1").fetchone()[0] or 0
    cobrado=conn.execute("SELECT SUM(monto_usd) FROM pagos").fetchone()[0] or 0
    benefs=conn.execute("SELECT * FROM beneficiarios WHERE activo=1 ORDER BY contrato").fetchall()
    conn.close()

    fname=f"Informe_Mensual_{date.today().strftime('%Y_%m')}.pdf"
    fpath=os.path.join(BASE_DIR,'reportes',fname)
    doc=SimpleDocTemplate(fpath,pagesize=A4,topMargin=2*cm,bottomMargin=2*cm,leftMargin=2*cm,rightMargin=2*cm)
    styles=getSampleStyleSheet()
    azul=colors.HexColor('#0C447C')
    
    story=[]
    story.append(Paragraph(f"🌱 {cfg.get('nombre_banco','BANCO COMUNITARIO')}", styles['Title']))
    story.append(Paragraph(f"Informe de Cartera - {date.today().strftime('%B %Y')}", styles['Normal']))
    story.append(HRFlowable(width="100%",thickness=2,color=azul,spaceAfter=12))
    
    resumen_data=[['Indicador','Valor'],['Total Beneficiarios',str(total)],['Cartera Total',f'${cartera:,.2f}'],['Monto Cobrado',f'${cobrado:,.2f}']]
    t=Table(resumen_data,colWidths=[10*cm,6*cm])
    t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),azul),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.5,colors.grey)]))
    story.append(t)
    
    doc.build(story)
    return send_file(fpath,as_attachment=True,download_name=fname,mimetype='application/pdf')

@app.route('/api/informe/pptx')
@login_required
def generar_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    conn=get_db()
    cfg={r['clave']:r['valor'] for r in conn.execute("SELECT clave,valor FROM config_sistema").fetchall()}
    total=conn.execute("SELECT COUNT(*) FROM beneficiarios WHERE activo=1").fetchone()[0]
    cartera=conn.execute("SELECT SUM(monto) FROM beneficiarios WHERE activo=1").fetchone()[0] or 0
    conn.close()
    
    prs=Presentation()
    slide=prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"🌱 {cfg.get('nombre_banco','BANCO COMUNITARIO')}"
    slide.placeholders[1].text = f"Informe de Gestión\n{date.today().strftime('%Y')}\nTotal Beneficiarios: {total}\nCartera: ${cartera:,.2f}"
    
    fname=f"Presentacion_{date.today().strftime('%Y_%m')}.pptx"
    fpath=os.path.join(BASE_DIR,'reportes',fname)
    prs.save(fpath)
    return send_file(fpath,as_attachment=True,download_name=fname,mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

# --- INICIO DE LA APP ---

if __name__ == '__main__':
    for d in ['datos','fotos','documentos','reportes','plantillas','backups']:
        os.makedirs(os.path.join(BASE_DIR,d),exist_ok=True)
    
    init_db()
    
    def abrir():
        import time
        time.sleep(1.5)
        webbrowser.open('http://localhost:5000/setup' if not is_setup() else 'http://localhost:5000')
    
    threading.Thread(target=abrir,daemon=True).start()
    
    print("\n"+"= "*52)
    print("  SemilleroComunal v2.3 (CACHE-FIX)")
    print("  Listo para usar")
    print("= "*52)
    print("  Navegador: http://localhost:5000")
    print("= "*52+"\n")
    
    app.run(debug=True, use_reloader=False, port=5000, host='127.0.0.1')